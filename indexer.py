import logging
import os
import argparse
from dotenv import load_dotenv

logger = logging.getLogger(__name__)
from llama_index.core import VectorStoreIndex, SimpleDirectoryReader, StorageContext
from llama_index.core.node_parser import SentenceSplitter
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.embeddings.huggingface import HuggingFaceEmbedding
from llama_index.readers.file import DocxReader, PandasExcelReader
from llama_index.vector_stores.qdrant import QdrantVectorStore
from qdrant_client import QdrantClient
from qdrant_client.models import Filter, FieldCondition, MatchValue, FilterSelector
from pathlib import Path
from PIL import Image
import pytesseract
import pypdf
from pdf2image import convert_from_path
import anthropic
import base64
import io

# Tesseract language pack — extend this string if more languages are needed
_OCR_LANG = "eng+ara+heb"
# PDFs with fewer characters than this per page are treated as scanned
_OCR_TEXT_THRESHOLD = 50

# Claude model used for vision OCR — far better than Tesseract for Arabic/Hebrew
_VISION_MODEL = "claude-haiku-4-5-20251001"
_OCR_PROMPT = (
    "Please extract ALL text from this image exactly as it appears. "
    "Preserve the original language (Arabic, Hebrew, English). "
    "Do not translate or summarize — just extract the raw text. "
    "If the image contains no readable text, return an empty response."
)

_anthropic_client = None


def _get_anthropic_client():
    global _anthropic_client
    if _anthropic_client is None:
        _anthropic_client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    return _anthropic_client


def _claude_ocr(image_b64: str, media_type: str) -> str:
    """Send a base64-encoded image to Claude and return the extracted text."""
    client = _get_anthropic_client()
    message = client.messages.create(
        model=_VISION_MODEL,
        max_tokens=4096,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": media_type,
                            "data": image_b64,
                        },
                    },
                    {"type": "text", "text": _OCR_PROMPT},
                ],
            }
        ],
    )
    if not message.content:
        return ""
    return message.content[0].text or ""


class SmartPDFReader(BaseReader):
    """Extracts text from PDFs; falls back to Claude Vision OCR for scanned/image-only pages."""

    def load_data(self, file, extra_info=None):
        path = Path(file)
        try:
            pages_text = []
            with open(path, "rb") as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    pages_text.append(page.extract_text() or "")

            avg_chars = sum(len(t) for t in pages_text) / max(len(pages_text), 1)
            if avg_chars < _OCR_TEXT_THRESHOLD:
                # Scanned/image-only PDF — rasterize each page and OCR via Claude,
                # which is dramatically better than Tesseract for Arabic/Hebrew.
                images = convert_from_path(str(path))
                pages_text = []
                for img in images:
                    buf = io.BytesIO()
                    img.save(buf, format="PNG")
                    page_b64 = base64.standard_b64encode(buf.getvalue()).decode("utf-8")
                    pages_text.append(_claude_ocr(page_b64, "image/png"))

            return [Document(text="\n\n".join(pages_text), metadata=extra_info or {})]
        except Exception as e:
            raise RuntimeError(f"Failed to read PDF '{path.name}': {e}") from e


class ClaudeVisionReader(BaseReader):
    """Uses Claude to extract text from images — much better than Tesseract for Arabic."""

    def load_data(self, file, extra_info=None):
        path = Path(file)
        try:
            with open(file, "rb") as f:
                image_data = base64.standard_b64encode(f.read()).decode("utf-8")

            ext = path.suffix.lower()
            media_type_map = {
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".gif": "image/gif",
                ".webp": "image/webp"
            }
            media_type = media_type_map.get(ext, "image/jpeg")

            extracted_text = _claude_ocr(image_data, media_type)
            return [Document(text=extracted_text, metadata=extra_info or {})]
        except Exception as e:
            raise RuntimeError(f"Failed to extract text from image '{path.name}': {e}") from e

class ImageOCRReader(BaseReader):
    """Extracts text from image files using Tesseract OCR."""

    def load_data(self, file, extra_info=None):
        path = Path(file)
        try:
            text = pytesseract.image_to_string(Image.open(file), lang=_OCR_LANG)
            return [Document(text=text, metadata=extra_info or {})]
        except Exception as e:
            raise RuntimeError(f"OCR failed for '{path.name}': {e}") from e


# Embedded images smaller than this (max dimension, px) are treated as decorative
# icons/logos and skipped, to avoid a wasted Claude call per bullet point.
_PPTX_MIN_IMAGE_DIM = 100


def _ocr_pptx_image(blob: bytes) -> str:
    """OCR a raw image blob from a PPTX slide via Claude. Returns '' if unusable."""
    try:
        img = Image.open(io.BytesIO(blob))
    except Exception:
        # Vector formats (EMF/WMF) and the like can't be rasterized here — skip.
        return ""
    if max(img.size) < _PPTX_MIN_IMAGE_DIM:
        return ""
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format="PNG")
    img_b64 = base64.standard_b64encode(buf.getvalue()).decode("utf-8")
    return _claude_ocr(img_b64, "image/png").strip()


class SmartPptxReader(BaseReader):
    """Extracts typed text, tables and speaker notes from PPTX, and OCRs embedded
    images via Claude — so mixed decks (typed Arabic + image-of-Arabic) index fully."""

    def load_data(self, file, extra_info=None):
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        path = Path(file)

        def walk(shapes, parts):
            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    walk(shape.shapes, parts)
                    continue
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        parts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append("\t".join(c.text.strip() for c in row.cells))
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        ocr = _ocr_pptx_image(shape.image.blob)
                        if ocr:
                            parts.append(ocr)
                    except Exception:
                        logger.exception("Image OCR failed on a slide in %s", path.name)

        try:
            prs = Presentation(str(file))
            slides_text = []
            for slide in prs.slides:
                parts = []
                walk(slide.shapes, parts)
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    note = slide.notes_slide.notes_text_frame.text.strip()
                    if note:
                        parts.append(note)
                if parts:
                    slides_text.append("\n".join(parts))
            return [Document(text="\n\n".join(slides_text), metadata=extra_info or {})]
        except Exception as e:
            raise RuntimeError(f"Failed to read PPTX '{path.name}': {e}") from e


SUPPORTED_EXTENSIONS = [
    ".pdf",
    ".txt",
    ".docx", ".pptx", ".xlsx", ".xls",
    ".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp",
]

_FILE_EXTRACTOR = {
    ".pdf":  SmartPDFReader(),
    ".docx": DocxReader(),
    ".pptx": SmartPptxReader(),
    ".xlsx": PandasExcelReader(),
    ".xls":  PandasExcelReader(),
    ".jpg":  ClaudeVisionReader(),   # ← Claude instead of Tesseract
    ".jpeg": ClaudeVisionReader(),   # ← Claude instead of Tesseract
    ".png":  ClaudeVisionReader(),   # ← Claude instead of Tesseract
    ".tiff": ImageOCRReader(),       # ← Keep Tesseract for TIFF
    ".tif":  ImageOCRReader(),
    ".bmp":  ImageOCRReader(),
}



load_dotenv()

COLLECTION_NAME = "rag_documents"

_client = None
_embed_model = None
_storage_context = None
_splitter = None


def _get_resources():
    global _client, _embed_model, _storage_context, _splitter
    if _client is None:
        try:
            _client = QdrantClient(url=os.getenv("QDRANT_URL", "localhost"), port=6333)
            _embed_model = HuggingFaceEmbedding(model_name="intfloat/multilingual-e5-large")
            vector_store = QdrantVectorStore(client=_client, collection_name=COLLECTION_NAME)
            _storage_context = StorageContext.from_defaults(vector_store=vector_store)
            _splitter = SentenceSplitter(chunk_size=512, chunk_overlap=50)
        except Exception as e:
            _client = None  # allow retry on next call
            raise RuntimeError(f"Failed to initialize indexer resources: {e}") from e
    return _client, _embed_model, _storage_context, _splitter


def _gather_supported_files(file_dir: str) -> list:
    """Recursively collect supported files, matching extensions case-insensitively.

    SimpleDirectoryReader's ``required_exts`` filter is case-sensitive, so a file
    named ``deck.PPTX`` (uppercase) would be silently skipped. We enumerate the
    files ourselves and pass them via ``input_files`` instead.
    """
    allowed = {e.lower() for e in SUPPORTED_EXTENSIONS}
    return sorted(
        str(p) for p in Path(file_dir).rglob("*")
        if p.is_file() and p.suffix.lower() in allowed
    )


def index_user_files(user_id: str, file_dir: str):
    if not os.path.isdir(file_dir) or not os.listdir(file_dir):
        raise ValueError(f"No files found in directory: {file_dir}")

    logger.info("Indexing files for user %s from directory %s", user_id, file_dir)
    client, embed_model, storage_context, splitter = _get_resources()

    input_files = _gather_supported_files(file_dir)
    if not input_files:
        raise ValueError(f"No supported files found in directory: {file_dir}")

    documents = SimpleDirectoryReader(
        input_files=input_files,
        file_extractor=_FILE_EXTRACTOR,
    ).load_data()
    for doc in documents:
        doc.metadata['user_id'] = user_id
    logger.info("Loaded %d documents for user %s", len(documents), user_id)

    VectorStoreIndex.from_documents(
        documents,
        storage_context=storage_context,
        embed_model=embed_model,
        transformations=[splitter],
        show_progress=True
    )
    logger.info("Indexed %d documents for user %s", len(documents), user_id)


def index_single_file(user_id: str, file_path: str):
    """Index exactly one file for a user, replacing any prior vectors for it.

    Used on upload/reindex so adding a file only embeds that file — unlike
    index_user_files, which re-reads the whole directory and would duplicate the
    vectors of every already-indexed file. Existing vectors for this file name
    are deleted first so re-indexing the same file never leaves stale duplicates.
    """
    path = Path(file_path)
    if not path.is_file():
        raise ValueError(f"File not found: {file_path}")
    if path.suffix.lower() not in {e.lower() for e in SUPPORTED_EXTENSIONS}:
        raise ValueError(f"Unsupported file type: {path.suffix}")

    logger.info("Indexing single file for user %s: %s", user_id, path.name)
    client, embed_model, storage_context, splitter = _get_resources()

    delete_file_vectors(user_id, path.name)

    documents = SimpleDirectoryReader(
        input_files=[str(path)],
        file_extractor=_FILE_EXTRACTOR,
    ).load_data()
    for doc in documents:
        doc.metadata['user_id'] = user_id
    logger.info("Loaded %d documents from %s for user %s", len(documents), path.name, user_id)

    VectorStoreIndex.from_documents(
        documents,
        storage_context=storage_context,
        embed_model=embed_model,
        transformations=[splitter],
        show_progress=True
    )
    logger.info("Indexed file %s for user %s", path.name, user_id)


def get_indexed_filenames(user_id: str) -> set:
    """Return the set of file names that have at least one indexed point for this user.

    Used to show indexing status in the UI. Returns an empty set if the collection
    does not exist yet or Qdrant is unreachable, so listing never fails on this.
    """
    try:
        client, _, _, _ = _get_resources()
    except Exception:
        logger.exception("Could not init resources while reading indexed filenames for %s", user_id)
        return set()

    user_filter = Filter(
        must=[FieldCondition(key="user_id", match=MatchValue(value=user_id))]
    )
    names = set()
    offset = None
    try:
        while True:
            points, offset = client.scroll(
                collection_name=COLLECTION_NAME,
                scroll_filter=user_filter,
                with_payload=["file_name"],
                with_vectors=False,
                limit=256,
                offset=offset,
            )
            for p in points:
                fn = (p.payload or {}).get("file_name")
                if fn:
                    names.add(fn)
            if offset is None:
                break
    except Exception:
        logger.exception("Failed to read indexed filenames for user %s", user_id)
        return set()
    return names


def delete_user_files(user_id: str):
    logger.info("Deleting indexed documents for user %s", user_id)
    client, _, _, _ = _get_resources()

    user_filter = Filter(
        must=[
            FieldCondition(
                key="user_id",
                match=MatchValue(value=user_id)
            )
        ]
    )

    client.delete(
        collection_name=COLLECTION_NAME,
        points_selector=FilterSelector(filter=user_filter)
    )
    logger.info("Deleted indexed documents for user %s", user_id)


def delete_file_vectors(user_id: str, file_name: str):
    """Delete only the vectors for a single file belonging to this user.

    Unlike delete_user_files (which drops every vector for the user), this
    filters on both user_id and file_name so deleting one file never touches
    the rest of the collection — and no re-index is required afterwards.
    """
    logger.info("Deleting indexed vectors for user %s, file %s", user_id, file_name)
    client, _, _, _ = _get_resources()

    file_filter = Filter(
        must=[
            FieldCondition(key="user_id", match=MatchValue(value=user_id)),
            FieldCondition(key="file_name", match=MatchValue(value=file_name)),
        ]
    )

    client.delete(
        collection_name=COLLECTION_NAME,
        points_selector=FilterSelector(filter=file_filter),
    )
    logger.info("Deleted vectors for user %s, file %s", user_id, file_name)


def main(args):
    index_user_files(user_id=args.user, file_dir=args.dir)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Indexer for document search')
    parser.add_argument('--user', type=str, required=True, help='Username for indexing and searching')
    parser.add_argument('--dir', type=str, required=True, help='Directory containing documents to index')
    args = parser.parse_args()
    main(args)
